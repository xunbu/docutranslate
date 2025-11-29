# SPDX-FileCopyrightText: 2025 QinHan
# SPDX-License-Identifier: MPL-2.0
import asyncio
from dataclasses import dataclass
from io import BytesIO
from typing import Self, Literal, List, Dict, Any, Tuple, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.text.text import _Paragraph, TextFrame
from pptx.oxml.ns import qn

from docutranslate.agents.segments_agent import SegmentsTranslateAgentConfig, SegmentsTranslateAgent
from docutranslate.ir.document import Document
from docutranslate.translator.ai_translator.base import AiTranslatorConfig, AiTranslator


# ---------------- 配置类 ----------------
@dataclass
class PPTXTranslatorConfig(AiTranslatorConfig):
    insert_mode: Literal["replace", "append", "prepend"] = "replace"
    separator: str = "\n"


# ---------------- 主类 ----------------
class PPTXTranslator(AiTranslator):
    """
    基于 python-pptx 的 .pptx 文件翻译器 (增强版)。

    改进特性：
    1. 深度遍历：支持母版、版式、备注页、以及隐藏在 AlternateContent (兼容性块) 中的文本。
    2. 公式保护：智能检测文本间的公式，防止翻译后文字错位。
    3. 样式保留：翻译后完全保留原有的中英文字体设置，不做强制覆盖。
    4. 布局自适应：防止翻译后文本溢出。
    """

    def __init__(self, config: PPTXTranslatorConfig):
        super().__init__(config=config)
        self.chunk_size = config.chunk_size
        self.translate_agent = None
        if not self.skip_translate:
            agent_config = SegmentsTranslateAgentConfig(
                custom_prompt=config.custom_prompt, to_lang=config.to_lang, base_url=config.base_url,
                api_key=config.api_key, model_id=config.model_id, temperature=config.temperature,
                thinking=config.thinking, concurrent=config.concurrent, timeout=config.timeout,
                logger=self.logger, glossary_dict=config.glossary_dict, retry=config.retry,
                system_proxy_enable=config.system_proxy_enable, force_json=config.force_json
            )
            self.translate_agent = SegmentsTranslateAgent(agent_config)
        self.insert_mode = config.insert_mode
        self.separator = config.separator

    # ---------------- 辅助函数：样式与字体 ----------------

    def _get_font_signature(self, run) -> Tuple:
        """获取 Run 的字体样式签名，用于合并判断。"""
        font = run.font
        color_key = None

        # 稳健的颜色获取逻辑
        if hasattr(font, 'color') and font.color:
            try:
                if font.color.type == MSO_COLOR_TYPE.RGB:
                    color_key = str(font.color.rgb)
                elif font.color.type == MSO_COLOR_TYPE.THEME:
                    color_key = f"THEME_{font.color.theme_color}_{font.color.brightness}"
            except AttributeError:
                pass

        return (
            font.name,
            font.size,
            font.bold,
            font.italic,
            font.underline,
            color_key
        )

    def _have_same_significant_styles(self, run1, run2) -> bool:
        """检查两个 Run 是否样式相同且在 XML 结构上紧邻（中间无公式）。"""
        if run1 is None or run2 is None:
            return False

        # 1. 检查视觉样式是否一致
        if self._get_font_signature(run1) != self._get_font_signature(run2):
            return False

        # 2. 检查 XML 邻接性
        # 如果 run1 和 run2 之间夹杂了 <m:oMath> (公式) 或其他标签，
        # 它们的 XML 索引将不连续。此时必须切分，否则回填时文字会跑到公式前面。
        try:
            r1_element = run1._r
            r2_element = run2._r
            parent = r1_element.getparent()

            # 只有当它们属于同一个父节点，且索引差为1时，才视为紧邻
            if parent == r2_element.getparent():
                index1 = parent.index(r1_element)
                index2 = parent.index(r2_element)
                if index2 != index1 + 1:
                    return False  # 中间有东西（如公式），禁止合并
        except Exception:
            # 如果底层操作失败，保守起见不合并
            return False

        return True

    # ---------------- 核心遍历逻辑 ----------------

    def _process_text_frame(self, text_frame: TextFrame, elements: List[Dict[str, Any]], texts: List[str]):
        """处理 TextFrame 中的所有段落"""
        for paragraph in text_frame.paragraphs:
            self._process_paragraph(paragraph, elements, texts)

    def _process_paragraph(self, paragraph: _Paragraph, elements: List[Dict[str, Any]], texts: List[str]):
        """处理单个段落，智能切分文本"""
        if not paragraph.runs:
            return

        current_runs = []

        def flush_segment():
            if not current_runs:
                return
            full_text = "".join(r.text for r in current_runs)
            # 只有非空文本才翻译
            if full_text.strip():
                elements.append({
                    "type": "text_runs",
                    "runs": list(current_runs),
                    "paragraph": paragraph,
                    "text_frame": paragraph._parent
                })
                texts.append(full_text)
            current_runs.clear()

        for run in paragraph.runs:
            # 这里的 run.text 只有纯文本，不包含公式内容
            if not run.text:
                continue

            last_run = current_runs[-1] if current_runs else None

            # 样式不同 或 物理位置不连续（中间有公式）则切分
            if last_run and not self._have_same_significant_styles(last_run, run):
                flush_segment()

            current_runs.append(run)

        flush_segment()

    def _process_shape(self, shape, elements: List[Dict[str, Any]], texts: List[str]):
        """递归处理常规形状"""
        # 1. 组合图形
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child_shape in shape.shapes:
                self._process_shape(child_shape, elements, texts)
            return

        # 2. 表格
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame") and cell.text_frame:
                        self._process_text_frame(cell.text_frame, elements, texts)
            return

        # 3. 常规文本框
        if shape.has_text_frame:
            try:
                self._process_text_frame(shape.text_frame, elements, texts)
            except Exception:
                pass

    def _scan_deep_xml_for_text(self, slide_element, elements: List[Dict[str, Any]], texts: List[str]):
        """
        [深度扫描] 直接遍历 XML 树，寻找标准 API 无法触及的文本。
        修复了 KeyError: 'mc' 问题。
        """
        # 定义 XML 命名空间 URI
        MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
        # 手动构建带命名空间的标签名，不依赖 qn()
        MC_ALT = f"{{{MC_NS}}}AlternateContent"
        MC_CHOICE = f"{{{MC_NS}}}Choice"

        # 对于 'p' (PresentationML) 命名空间，python-pptx 支持 qn，可以继续使用
        P_SP = qn('p:sp')
        P_TXBODY = qn('p:txBody')

        # 查找所有 AlternateContent 块
        for alt_content in slide_element.iter(MC_ALT):
            # 找到 Choice 分支
            choice = alt_content.find(MC_CHOICE)
            if choice is None:
                continue

            # 在 Choice 内部寻找形状 (p:sp)
            for sp in choice.iter(P_SP):
                # 寻找 p:txBody (文本主体)
                txBody = sp.find(P_TXBODY)
                if txBody is not None:
                    try:
                        # 手动构建 TextFrame 对象
                        # 这里的 parent 设为 None 在读取/写入 text 属性时通常是安全的
                        tf = TextFrame(txBody, None)
                        self._process_text_frame(tf, elements, texts)
                    except Exception as e:
                        self.logger.warning(f"处理深度 XML 文本框时出错: {e}")

    def _scan_presentation_content(self, prs: Presentation, elements: List[Dict[str, Any]], texts: List[str]):
        """全量扫描 PPT 内容"""

        # 辅助内部函数：扫描单个“幻灯片类”对象
        def scan_slide_object(slide_obj):
            # 1. 常规 API 遍历 (处理普通文本、表格、组合)
            for shape in slide_obj.shapes:
                self._process_shape(shape, elements, texts)

            # 2. 深度 XML 遍历 (处理 AlternateContent/公式文本)
            self._scan_deep_xml_for_text(slide_obj.element, elements, texts)

        # 1. 遍历普通幻灯片 (Slides)
        for slide in prs.slides:
            scan_slide_object(slide)
            # 备注页
            if slide.has_notes_slide:
                notes = slide.notes_slide
                if notes.notes_text_frame:
                    self._process_text_frame(notes.notes_text_frame, elements, texts)

        # 2. 遍历母版 (Slide Masters)
        for master in prs.slide_masters:
            scan_slide_object(master)

            # 3. 遍历版式 (Layouts)
            for layout in master.slide_layouts:
                scan_slide_object(layout)

    # ---------------- 翻译前后处理 ----------------

    def _pre_translate(self, document: Document) -> Tuple[Presentation, List[Dict[str, Any]], List[str]]:
        """解析 PPT 文件"""
        prs = Presentation(BytesIO(document.content))
        elements, texts = [], []

        self._scan_presentation_content(prs, elements, texts)
        self.logger.info(f"共提取了 {len(texts)} 个文本片段 (包含隐藏的公式文本)。")
        return prs, elements, texts

    def _apply_translation(self, element_info: Dict[str, Any], final_text: str):
        """回填翻译，精细控制样式"""
        runs = element_info["runs"]
        if not runs:
            return

        original_text = "".join(r.text for r in runs)

        text_to_set = final_text
        if self.insert_mode == "append":
            text_to_set = original_text + self.separator + final_text
        elif self.insert_mode == "prepend":
            text_to_set = final_text + self.separator + original_text

        # --- 回填策略 ---
        primary_run = runs[0]

        try:
            # 1. 写入文本 (python-pptx 会自动保留原有的 rPr 属性，即保留默认字体)
            primary_run.text = text_to_set

            # 2. (已移除字体强制设置逻辑，以保留 PPT 原样)

            # 3. 处理溢出
            text_frame = element_info.get("text_frame")
            if text_frame and hasattr(text_frame, 'auto_size'):
                if text_frame.auto_size == MSO_AUTO_SIZE.NONE:
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        except Exception as e:
            self.logger.warning(f"应用翻译到 Run 时出错: {e}")
            return

        # 清空后续 run (模拟合并效果)
        for i in range(1, len(runs)):
            runs[i].text = ""

    def _after_translate(self, prs: Presentation, elements: List[Dict[str, Any]], translated: List[str],
                         originals: List[str]) -> bytes:
        """保存结果"""
        if len(elements) != len(translated):
            min_len = min(len(elements), len(translated))
            elements = elements[:min_len]
            translated = translated[:min_len]

        for info, trans in zip(elements, translated):
            self._apply_translation(info, trans)

        output_stream = BytesIO()
        prs.save(output_stream)
        return output_stream.getvalue()

    # ---------------- 接口实现 ----------------

    def translate(self, document: Document) -> Self:
        prs, elements, originals = self._pre_translate(document)
        if not originals:
            self.logger.info("未找到可翻译文本。")
            document.content = self._after_translate(prs, elements, [], [])
            return self

        if self.glossary_agent:
            self.glossary_dict_gen = self.glossary_agent.send_segments(originals, self.chunk_size)
            if self.translate_agent:
                self.translate_agent.update_glossary_dict(self.glossary_dict_gen)

        translated = self.translate_agent.send_segments(originals,
                                                        self.chunk_size) if self.translate_agent else originals
        document.content = self._after_translate(prs, elements, translated, originals)
        return self

    async def translate_async(self, document: Document) -> Self:
        prs, elements, originals = await asyncio.to_thread(self._pre_translate, document)
        if not originals:
            self.logger.info("未找到可翻译文本。")
            document.content = await asyncio.to_thread(self._after_translate, prs, elements, [], [])
            return self

        if self.glossary_agent:
            self.glossary_dict_gen = await self.glossary_agent.send_segments_async(originals, self.chunk_size)
            if self.translate_agent:
                self.translate_agent.update_glossary_dict(self.glossary_dict_gen)

        translated = await self.translate_agent.send_segments_async(originals,
                                                                    self.chunk_size) if self.translate_agent else originals
        document.content = await asyncio.to_thread(self._after_translate, prs, elements, translated, originals)
        return self